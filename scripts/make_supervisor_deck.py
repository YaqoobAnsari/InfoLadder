"""Supervisor-meeting deck generator: one slide spec -> PPTX + PDF.

Regenerate:  $TOPOSPEC_PYTHON scripts/make_supervisor_deck.py
Outputs: results/presentations/2026-07-21_supervisor_meeting_01.{pptx,pdf}

Every number in the spec traces to a repo artifact (registry runs, corpus trees,
batch reports, the audit doc); nothing is hand-invented. Slides are deliberately
low-density (max ~7 bullets); depth comes from slide count, not clutter.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "results" / "presentations"
STEM = "2026-07-21_supervisor_meeting_01"

# palette (matches topospec.reporting.style / validated dataviz set)
INK = "#0b0b0b"
SUB = "#52514e"
ACC = "#2a78d6"
BG = "#fcfcfb"
ACC2 = "#008300"
WARN = "#e34948"

IMG = {
    "t0": ROOT / "results/corpora/prelim_rasters/FF_part_1upE/overlays/T0.png",
    "t1": ROOT / "results/corpora/prelim_rasters/FF_part_1upE/overlays/T1.png",
    "t2": ROOT / "results/corpora/prelim_rasters/FF_part_1upE/overlays/T2.png",
    "t3": ROOT / "results/corpora/prelim_rasters/FF_part_1upE/overlays/T3.png",
    "src": ROOT / "results/corpora/prelim_rasters/FF_part_1upE/source.png",
    "calib": ROOT / "results/reports/a0_calibration/figures/calibration_surface.png",
    "controls": ROOT / "results/reports/a0_calibration/figures/controls.png",
    "contact": ROOT / "data/derived/msd_render/contact_sheet.png",
    "arcreal": ROOT / "data/derived/msd_render/door_arc_vs_real.png",
    "crops": ROOT / "data/derived/msd_door_train/sample_crops.png",
    "evalsynth": ROOT / "runs/door_retrain_v1/eval_synth.png",
}

S = []  # the deck spec


def title(t, sub, meta):
    S.append({"kind": "title", "title": t, "sub": sub, "meta": meta})


def section(n, t, blurb=""):
    S.append({"kind": "section", "n": n, "title": t, "blurb": blurb})


def bullets(t, items, note=""):
    S.append({"kind": "bullets", "title": t, "items": items, "note": note})


def table(t, headers, rows, note=""):
    S.append({"kind": "table", "title": t, "headers": headers, "rows": rows, "note": note})


def image(t, key, caption=""):
    S.append({"kind": "image", "title": t, "img": IMG[key], "caption": caption})


def imagebul(t, key, items, caption=""):
    S.append({"kind": "imagebul", "title": t, "img": IMG[key], "items": items, "caption": caption})


def statement(t, lines):
    S.append({"kind": "statement", "title": t, "lines": lines})


# ============================== CONTENT ======================================

title(
    "InfoLadder",
    "How Much Structure Does Physical Inference Need?\nMeasuring Usable Information Across a Floorplan Representation Spectrum",
    ["Weekly supervision meeting 1  ·  2026-07-21",
     "Yaqoob Ansari  ·  TopoField thread T4  ·  target: ICLR (mid-September 2026)",
     "github.com/YaqoobAnsari/InfoLadder"],
)

statement("The thesis", [
    "A floorplan is not a picture — it is a deliberately engineered encoding.",
    "Architects put circulation, access control, zoning and egress logic INTO the drawing by convention.",
    "That information is present but locked in pixels: not bioavailable to models.",
    "Graph representations of increasing richness progressively unlock it.",
    "We measure — for the first time under a formal protocol — how much each tier unlocks,",
    "and what each tier costs to build.",
])

# ---------------- A: context ----------------
section("A", "Context, question, and research gap")

bullets("The problem", [
    ("Building-graph representations range from bare room connectivity to typed, directed, hierarchical graphs.", 0),
    ("Nobody knows which structural content downstream inference actually needs.", 0),
    ("Consequences today:", 0),
    ("annotation effort is spent blind — richness is added by habit, not by measured value;", 1),
    ("architectures are chosen by convention (“use a GNN”) without knowing what the input carries;", 1),
    ("industry tiers information need (ISO 7817-1:2024) with NO quantitative selection method.", 1),
])

bullets("The formal research question", [
    ("For a nested ladder of representations T0 ⊆ T1 ⊆ … ⊆ T5 of the same floorplan,", 0),
    ("measure usable information I_V(T_k → Y) — what a bounded predictor family V can extract about target Y —", 0),
    ("as a curve over tiers, per task, per predictor capacity.", 0),
    ("Derived quantities:", 0),
    ("annotation gain: nats of usable information per human annotation hour, per tier increment (first in the literature);", 1),
    ("capacity compensation: can bigger models on poor representations match small models on rich ones, at what sample cost;", 1),
    ("three currencies — structure (hours), capacity (compute), samples (data) — priced against each other.", 1),
])

bullets("Why this matters now (verified evidence)", [
    ("ISO 7817-1:2024 “Level of Information Need” standardizes exactly this richness tiering — with no measurement method behind it.", 0),
    ("Scan-to-BIM semantic enrichment costs 200–400 manual hours per project — the cost axis we price.", 0),
    ("A 2026 publication wave (FloorplanQA/ICML26, FMLM/CVPR26 Highlight, floorplan tokenization, SVG-decomposition studies) shows the field circling the question without the measurement.", 0),
    ("The measurement tool itself (predictive V-information, ICLR 2020) is mature, trusted, and has never been pointed at input representation tiers.", 0),
])

table("Research landscape (adversarially verified, July 2026)",
      ["Work", "What it does", "What it lacks vs. ours"],
      [
          ["V-information line (Xu ICLR'20; Ethayarajh ICML'22)", "usable info of embeddings / dataset difficulty", "never input tiers, cost, or graphs"],
          ["Almudévar & Ortega '26", "V-info to compare learned representations", "hidden activations, not input formats"],
          ["Clio (RSS'24, MIT-SPARK)", "IB-selects scene-graph granularity for a task", "online compression, no measurement/cost/capacity"],
          ["KG-construction benchmark '26", "construction choice scored by downstream GNN", "accuracy-only, unordered variants, no cost"],
          ["FloorplanQA (ICML'26)", "LLM accuracy across floorplan text formats", "prompting eval; no nesting, probing, or cost"],
          ["Homophily / label-informativeness (Platonov '23)", "when structure helps, on fixed graphs", "binary use-vs-ignore; not which tier to build"],
      ],
      note="Full threat analysis with sources: docs/audits/project_audit_2026-07-16.md")

bullets("The gap, precisely", [
    ("“Representation matters” is folklore — ablations, domain observations, and 2026 benchmarks all attest to it.", 0),
    ("What has never existed, in any domain, is the combination:", 0),
    ("a NESTED refinement spectrum (each tier provably recoverable from the next by forgetting);", 1),
    ("measured under a formal usable-information protocol with probe-CAPACITY sweeps;", 1),
    ("tied to the ANNOTATION COST of each refinement;", 1),
    ("on an instrument CALIBRATED against planted ground truth before any real-data claim.", 1),
    ("No published relative combines any two of these four ingredients.", 0),
])

bullets("What we claim — and what we deliberately do not", [
    ("Claim: the six-tier spectrum formalization with cost semantics (C1).", 0),
    ("Claim: the first usable-information measurement over input representation tiers, as a calibrated instrument (C2).", 0),
    ("Claim: annotation gain in nats/hour (C3); information-richness curves + capacity/scale analyses (C4); domain generality check (C5).", 0),
    ("NOT claimed: a new probing method (we apply Xu et al.), a new GNN, or deployed physical-field inference.", 0),
    ("Wording discipline: the novelty is the PROTOCOL — “first to study representation effects on floorplans” is no longer defensible and is never said.", 0),
])

bullets("Measurement framework in one slide", [
    ("Predictive V-information (Xu et al., ICLR 2020): H_V(Y|X) = best expected log-loss achievable by predictor family V; I_V(X→Y) = H_V(Y) − H_V(Y|X).", 0),
    ("Key property: unlike Shannon MI, PROCESSING CAN INCREASE usable information — annotation is computation done in advance.", 0),
    ("Because tiers are nested, every tier is a coarsening of T5: the study is V-information under structured coarsening.", 0),
    ("MDL prequential codelength (Voita & Titov 2020) reported alongside — the probe-effort-sensitive, optimization-robust complement.", 0),
    ("All estimates are lower bounds achieved by optimization → best-of-3 restarts × 3 seeds, never single probes.", 0),
])

# ---------------- B: the spectrum ----------------
section("B", "The six-tier representation spectrum",
        "Every tier derives from the raster floorplan through one real parsing pipeline; downgrades are deterministic forgetting maps.")

table("The ladder at a glance",
      ["Tier", "Adds", "Derivation", "Cost"],
      [
          ["T0  skeleton", "spaces (rooms, corridor hubs, stairs/elevators) + undirected connectivity; centroids kept", "automatic (forget of T2)", "free"],
          ["T1  named spaces", "node kinds (room/corridor/transition) + recognized text labels", "automatic (CRAFT text)", "free"],
          ["T2  doored structure", "door NODES with subtypes (room–corridor, room–room, corridor–corridor, exit)", "automatic (door detector)", "free"],
          ["T3  measured geometry", "per-node measures: area, eq. radius, inradius, subnode count, door count", "automatic (pipeline stats)", "free"],
          ["T4  access control", "direction/restriction on connections (one-way, staff-only)", "MANUAL annotation, timed", "≈0.5–1 h/bldg"],
          ["T5  organization", "zone/wing containment hierarchy + outdoor marking", "MANUAL annotation, timed", "≈1–2 h/bldg"],
      ],
      note="Strict refinement: forget(T_{k+1}) = T_k, tested property. Thermal/crowding annotations are TARGETS, never representation content (leakage).")

imagebul("T0 — untyped skeleton", "t0", [
    ("What: every space is a node with only its position; edges mean “connected through an opening”.", 0),
    ("Why this floor: the minimum a model could receive — pure topology + layout.", 0),
    ("Every question answerable here is free; everything above is measured against this floor.", 0),
], caption="Institutional test building (FF sheet), T0: 42 spaces, untyped, connectivity only.")

imagebul("T1 — named spaces", "t1", [
    ("Adds: node kinds (room / corridor / stairs / elevator) + the drawing's text (room numbers, ‘hall’).", 0),
    ("Why: the first semantic increment — what each thing IS; still free (text is read from the sheet).", 0),
    ("Semantic key on our sheets: numbers = rooms, ‘hall’ = corridors, ‘NA’ = outdoor (excluded).", 0),
], caption="T1: same structure, now typed and labeled; corridors are the hall-text identities.")

imagebul("T2 — doored structure", "t2", [
    ("Adds: doors as first-class typed nodes (room–corridor, room–room, corridor–corridor, exit).", 0),
    ("Why: connection TYPE is physical structure — a wall shares heat, a door passes people.", 0),
    ("This is where circulation logic becomes explicit; detector-derived on rasters, annotation-derived on ArchCAD.", 0),
], caption="T2: 42 spaces + 38 typed door nodes + 79 edges; doors sit on the drawn openings.")

imagebul("T3 — measured geometry", "t3", [
    ("Adds: numeric measures per space — area, equivalent radius, inradius, sub-node count, door count.", 0),
    ("Why: physical capacity and shape complexity, surfaced as features; still fully automatic.", 0),
    ("T3 closes the FREE tiers: everything above requires human hours.", 0),
], caption="T3: node size now encodes measured area; all measures carried as attributes.")

bullets("T4 + T5 — the paid tiers (built only if the stage gate opens)", [
    ("T4 access control: direction/restriction per connection (one-way exits, staff-only doors) — the egress-relevant increment; ~0.5–1 h/building, timed.", 0),
    ("T5 organization: zone/wing containment forest + outdoor regions — the facility-manager's mental model; ~1–2 h/building, timed.", 0),
    ("Timing sheets make these the first measured c(k→k+1) figures — the denominator of annotation gain (C3).", 0),
    ("An annotation tool over the pipeline's pre-fill is designed but deliberately deferred (future TODO).", 0),
])

bullets("Why nesting is load-bearing (not a design nicety)", [
    ("Each tier is recoverable from the one above by a deterministic forgetting map φ — tested: determinism, idempotence, chain composition.", 0),
    ("This makes the ladder a measurement SCALE: differences between tiers are attributable to the added content, nothing else.", 0),
    ("Proposition 1: if the probe family is closed under composition with φ, I_V is monotone in tier — closure failure makes non-monotonicity (S9) a testable prediction, not an anecdote.", 0),
    ("Schema v2 enforces tier-licensed content (doors only at T2+, measures at T3+, direction at T4+, hierarchy at T5) — 113 automated tests.", 0),
])

bullets("Stage-gated execution (risk control)", [
    ("Phase 1 (now): build T0–T3 at corpus scale — thousands of accurate, visually inspectable graphs.", 0),
    ("Phase 2: run the full probing grid on the four FREE tiers.", 0),
    ("GATE: only an established, statistically supported information ordering across T0–T3 unlocks the manual T4/T5 investment.", 0),
    ("Rationale: never spend 40 annotation hours before the instrument demonstrates the effect on free tiers.", 0),
])

# ---------------- C: pipeline ----------------
section("C", "Raster-first construction pipeline",
        "One real pipeline builds every tier from the pixel input — the thesis is honored in the construction, not just the framing.")

bullets("Pipeline architecture", [
    ("Engine: Tesseract2 (our own SIGSPATIAL lineage) — CRAFT text detection → semantic interpreter → flood-fill space segmentation → Faster R-CNN door detection → navigable graph.", 0),
    ("Tier factory (this project): contracts the navigation graph into tier graphs —", 0),
    ("room sub-nodes → parent rooms; corridor waypoint meshes → hall-text identities (spatial territory assignment);", 1),
    ("recognized text joined deterministically from the interpreter sidecar (rooms=numbers, hall, NA→outside);", 1),
    ("typed doors, stairs/elevators as transitions; measures from pipeline statistics → validated T3 → forget() → T2/T1/T0.", 1),
    ("Fallback lane: classical morphology (watershed on wall masks) — kept as QA cross-check and for annotation-free sources.", 0),
])

image("Flagship result — the institutional test building at T2", "t2",
      caption="FF sheet: 42 spaces, 38 typed doors, 79 edges. Every numbered room is one node; doors sit on drawn openings; "
              "each hall text is a corridor identity serving its own territory. This graph passed the owner's manual inspection.")

bullets("Validating the pipeline against ground truth (the MSD loop)", [
    ("MSD (ECCV'24) ships true room polygons, types, and access graphs for 4,167 usable plans.", 0),
    ("We RENDER its geometry to clean rasters (walls + room-type text + door symbols), run the FULL raster pipeline, and score the output against the shipped truth.", 0),
    ("This turns ‘is your graph even right?’ into a table: room detection F1, type accuracy, adjacency agreement — at thousands-of-plans scale.", 0),
    ("Renders are deterministic (seeded per plan) and auditable; the room-type→vocabulary remap is recorded per plan.", 0),
])

image("MSD ground-truth renders (validation corpus)", "contact",
      caption="20-plan render sample: walls from true polygons, room-type text for the text stage, door symbols at true openings. "
              "4,167 plans rendered at 30–60 px/m.")

table("Pipeline accuracy vs. MSD ground truth (first 119 scored plans)",
      ["Metric", "Value", "Interpretation"],
      [
          ["Room detection precision / recall / F1", "0.996 / 0.977 / 0.986", "space segmentation is essentially solved"],
          ["Room-type accuracy (text → label)", "0.823 batch; 16/16 after vocabulary fixes", "reading is clean; residual = vocabulary mapping"],
          ["Adjacency F1 (door-contracted)", "0.125", "the door bottleneck, quantified (below)"],
          ["Door-symbol detection rate", "34.9% (3,687 rendered arcs)", "style gap vs. detector's training distribution"],
          ["Isolated rooms after graphing", "0", "connectivity is preserved even pre-door-fix"],
      ],
      note="Full-corpus re-scoring in flight (2,883 graphs built and growing; failure accounting reconciles every missing output).")

imagebul("The door bottleneck, diagnosed", "arcreal", [
    ("The detector was trained on bold double-swing symbols seated in thick hatched walls (top row).", 0),
    ("Naive synthetic arcs (bottom) are thin, isolated, sometimes overlapping → 2/19 detected on the gate plan.", 0),
    ("Decision: retrain the detector on synthetic doors — unlimited exact labels are generatable from ground truth.", 0),
], caption="Real training-distribution doors vs. first-generation synthetic arcs — the style/context gap in one figure.")

image("Door training set: 12,290 exact boxes across 400 plans", "crops",
      caption="Bounding boxes emitted analytically at render time (nothing detected, nothing hand-labeled). Domain randomization: "
              "wall thickness/line-style/hatching, stroke weight, single/double swing, hinge side, jambs, text presence, rotations. "
              "Plan-level 15% holdout; 100% ink-coverage self-check.")

bullets("Detector fine-tune: first iteration results", [
    ("Held-out synthetic recall: 0.048 (original detector) → 0.523 after 18 minutes of fine-tuning (precision 0.817) — an 11× gain, gate is ≥0.80.", 0),
    ("On the REAL sheets: +38% more detections while preserving 84.6% of the original detections' positions —", 0),
    ("pending manual QA: extra doors real (gate was mis-specified) vs. hallucinated (needs anchoring).", 1),
    ("v2 planned: longer training, per-subtype labels, real-sheet verdict folded in.", 0),
    ("Deployment discipline: original checkpoint preserved; nothing deploys until both gates pass AND the running corpus finishes (consistency).", 0),
])

bullets("Robustness engineering (what broke and what we did)", [
    ("Segmentation-palette crash: plans without corridors/balconies crashed the door classifier (~36% of MSD) — root-caused, minimally patched, verified; patch archived for upstreaming.", 0),
    ("Every missing output is ACCOUNTED: failure fingerprinting buckets each absent graph by its crash stage — nothing is silently lost.", 0),
    ("Batches are resumable and idempotent; renders are seeded-deterministic; the failed 785 plans are being swept right now with the patch live (0 new failures so far).", 0),
])

# ---------------- D: data ----------------
section("D", "Corpora — what we have, quantified")

table("Corpus overview",
      ["Corpus", "Scale", "Annotations", "Doors", "Role"],
      [
          ["Prelim pool (ours)", "21 sheets: 1 institutional building (FF×8, LF×2, SF×1 + variants) + 10 residential", "text labels (numbers / hall / NA)", "drawn symbols (detector)", "pipeline testbed; Gate-b gold candidate"],
          ["MSD (ECCV'24)", "5,372 public; 4,167 usable; ≈140k spaces", "room polygons, 9 types, access graphs, zones", "edge attribute only (no geometry)", "GROUND TRUTH validation + residential scale"],
          ["FloorPlanCAD (ICCV'21)", "11,602 SVGs (6,965/3,827/810)", "35 classes, line-grained, per-primitive", "drawn symbols, 6 classes", "parser template; textless-lane source"],
          ["ArchCAD-400K (NeurIPS'25)", "41,097 slices, 14 m × 14 m", "numeric per-primitive semantics", "GROUND-TRUTH typed (4 classes + instances)", "institutional regime (S4); survey-gated"],
          ["Code graphs (planned)", "1 established defect dataset", "AST → +dataflow → +call ladder", "n/a", "generality check C5, hard-capped"],
      ])

bullets("Prelim pool — the owner-verified anchor (quantified)", [
    ("21 sheets; the institutional building spans FF (8 sheets), LF (2), SF (1) with text-annotation variants deduplicated.", 0),
    ("Tier graphs at T3: 359 spaces, 368 typed doors, 733 edges; 17.1 spaces/sheet average; 100% of spaces carry labels.", 0),
    ("Text key: rooms = numbers (e.g., 1001–1040), corridors = ‘hall’, outdoor = ‘NA’, stairs/elevator recognized.", 0),
    ("Quality: manually inspected by the building's owner; corridor territories and labels verified against reality.", 0),
    ("Residential file_N renders excluded from measurement (furniture-at-wall-thickness defeats morphology; documented).", 0),
])

bullets("MSD — Modified Swiss Dwellings (validation + residential scale)", [
    ("5,372 public multi-unit plans; ONLY the 4,167 train plans ship geometry (test withholds it) — corpus arithmetic corrected early.", 0),
    ("≈140k spaces; 33.6 spaces/plan average (measured on the 200-plan subset: 6,720 spaces).", 0),
    ("Room types (9 area classes), access edges typed door/passage/entrance; NO access direction; NO door geometry (approximated at nearest boundary points when rendering).", 0),
    ("Leakage found & neutralized: shipped zones are a deterministic function of room type → zone targets use APARTMENT grouping (access components cut at entrance doors), which is genuinely tier-exclusive.", 0),
    ("License CC BY 4.0; 2,883 plans already carried through the full pipeline to tier graphs.", 0),
])

bullets("FloorPlanCAD + ArchCAD-400K — the institutional axis", [
    ("FloorPlanCAD: 11,602 annotated SVGs incl. schools/hospitals/malls; 35 line-grained classes; doors drawn (6 classes).", 0),
    ("Caveat found by exhaustive search: every sheet is a 140×140-unit CROP; zero self-contained plans in the val split → role restricted to parser template + textless-lane source.", 1),
    ("ArchCAD-400K (gated access approved): 41,097 slices of professional drawings; walls/glass/columns/stairs/elevators labeled;", 0),
    ("doors are GROUND-TRUTH typed instances → the door tier needs no detector on this corpus;", 1),
    ("no room polygons and no text → rooms derived by our watershed lane; an ENCLOSURE SURVEY across all slices gates admission (crops again).", 1),
    ("Licenses: CC BY-NC 4.0 (both), ArchCAD additionally gated non-commercial; research-use, cited, never redistributed.", 0),
])

bullets("Data risk register (found early, neutralized)", [
    ("Zone-label determinism (MSD): would leak the T5 target through the T1 room label → apartment-mode zones adopted as default.", 0),
    ("Base-rate channel (synthetic): uneven per-building label rates let probes read building identity, not labels — caught by the calibration instrument itself, closed by balanced design.", 0),
    ("Crop fragmentation (FloorPlanCAD, ArchCAD): admission gates by enclosure, usable fractions reported, never assumed.", 0),
    ("Detector style-transfer (doors): quantified (34.9%), retrain program running, ArchCAD path is annotation-based and unaffected.", 0),
])

# ---------------- E: instrument ----------------
section("E", "The calibrated measurement instrument")

bullets("Why the instrument must be calibrated first", [
    ("Every I_V estimate is a lower bound achieved by optimization — an uncalibrated estimator is an article of faith.", 0),
    ("Design: synthetic buildings where each target is PLANTED at a known tier by construction —", 0),
    ("degree (T0), semantic label (T1), door subtype (T2), numeric attribute (T3), direction (T4), zone secret (T5).", 1),
    ("PASS iff the measured surface jumps exactly at each planted tier, is flat elsewhere, and shuffled controls extract nothing.", 0),
    ("This doubles as the paper's Figure 2 — the instrument's license to make real-data claims.", 0),
])

image("Calibration result: PASS — 71/71 checks (V1–V5) + 19/19 (V6, GPU)", "calib",
      caption="Each panel is one planted target; the dashed line marks where the answer was hidden. Every readable family jumps exactly "
              "at the planted tier and is flat below it; the V0 parameter-free readout at T5 extracts 0.73 of the 0.73 available nats — "
              "richness converts learning into lookup.")

bullets("The instrument earned trust by catching real bugs", [
    ("Base-rate channel: shuffled-label controls extracted +0.115 nats via building-identifying features — a 2/53-check FAIL that exposed a design flaw, fixed by balanced planting.", 0),
    ("Degenerate MDL blocks: single-class training prefixes crashed the linear family — now degrades to the marginal predictor (regression-tested).", 0),
    ("Mean-aggregation blindness: mean-pooled GNNs provably cannot count — switched to sum aggregation after the degree control failed.", 0),
    ("Each catch is documented as a decision record; failed runs stay in the registry — nothing is retried-until-green silently.", 0),
])

table("The probe ladder (the capacity sweep V)",
      ["Rung", "Family", "Question it answers"],
      [
          ["V0", "parameter-free readout (0 params)", "is the answer literally written in the representation?"],
          ["V1", "predict-from-prior", "the floor: what does the label marginal give?"],
          ["V2 / V3", "linear / linear + spectral encodings", "is the information on the surface of node features?"],
          ["V4 / V5", "1- / 2-layer GNNs (sum aggregation)", "is it extractable with local computation?"],
          ["V6", "4-layer GraphGPS-style transformer (≤ 2M params, GPU)", "is it extractable at all (bounded ceiling)?"],
          ["V7 (planned)", "frozen small LM over serialized graphs", "pre-empts ‘why not an LLM?’ (appendix tier)"],
      ],
      note="Fixed architectures across tiers (fairness); every curve sandwiched between shuffled-label floors and oracle skylines.")

bullets("Evaluation metric suite (revised after methods review)", [
    ("Core: I_V per (tier, target, family, seed) + prequential MDL codelength — both in nats, both released raw.", 0),
    ("Tier increments measured as CONDITIONAL V-information I_V(ΔT → Y | T_k) — the established ‘information beyond a baseline’ construction, not noisy differencing.", 0),
    ("Adopted additions (each neutralizes a known objection at marginal cost):", 0),
    ("pointwise V-information + dataset-artifact audit (per-instance difficulty; Data-Checklist lineage, COLM 2024);", 1),
    ("loss-data-curve summaries from the MDL blocks we already compute (sample-efficiency view).", 1),
    ("Considered and rejected with reasons: PID, graph information bottleneck, neural MI estimators, effective information, intrinsic dimension, concept-erasure probing.", 0),
])

bullets("Statistical protocol (pre-registered before the grid runs)", [
    ("Building-level splits, never floor-level; split assignment hashed and frozen.", 0),
    ("3 seeds × best-of-3 restarts; 1000-resample cluster bootstrap (clusters = buildings) for every CI.", 0),
    ("Paired Wilcoxon across buildings for tier comparisons, Holm-corrected within each tier family.", 0),
    ("Hard validity gates: shuffled controls must extract nothing (one-sided); oracle skylines must dominate; calibration must be PASS.", 0),
    ("Negative and flat results are published with equal prominence — committed in writing.", 0),
])

bullets("Methodological guard: separating structure from annotation content", [
    ("Confound: adding door nodes (T2) or hierarchy (T5) changes graph TOPOLOGY — commute times and effective resistance shift — so gains could be partly ‘free structure’.", 0),
    ("Controls added to the grid design:", 0),
    ("free-structure baselines: virtual-node / Louvain-clustering hierarchies vs. the manual T5; untyped added nodes vs. typed doors;", 1),
    ("per-tier effective-resistance deltas reported; an MLP ablation isolates feature- from structure-carried information.", 1),
    ("Without these, nats-per-annotation-hour is attackable; with them, it is the paper's most defensible number.", 0),
])

# ---------------- F: targets & plan ----------------
section("F", "Targets, applications, hypotheses, and plan")

table("Target suite with 2D-validity verdicts",
      ["Target", "Definition", "Verdict"],
      [
          ["Y_pde", "steady-state diffusion field solved exactly on the true 2D geometry", "KEEP — reframed as an analytic geometric functional; exact by construction, immune to simulator objections"],
          ["Y_zone", "plan-level thermal zoning (ASHRAE-grounded archetypes, assumptions disclosed)", "KEEP — zoning is a plan-level design activity"],
          ["Y_rank", "pairwise room ordering on the diffusion field", "KEEP, DEMOTED — renamed away from ‘temperature prediction’ (field data refute plan-only prediction of real temperatures)"],
          ["Y_egress", "evacuation bottleneck membership (+ betweenness reference)", "PROMOTED to co-headline — fully 2D-native, high-impact"],
          ["Y_syntax (proposed)", "space-syntax integration values", "ADD (pending decision) — zero annotation cost, established field"],
      ],
      note="Paper gains a ‘2D validity envelope’: stack effect, solar orientation, inter-floor transfer are explicitly out of scope. "
           "This concedes exactly what should be conceded — and the study gets stronger.")

bullets("Downstream applications: who needs this measurement", [
    ("WELCOMES (anchor the paper here):", 0),
    ("scan-to-BIM semantic enrichment — 200–400 manual h/project is precisely the cost we price;", 1),
    ("ISO 7817-1 Level-of-Information-Need decisions — standardized tiers, no quantitative selection method;", 1),
    ("LLM building QA / reasoning — input-format choice is live and unmeasured (FloorplanQA line).", 1),
    ("INDIFFERENT: indoor localization, generic navigation (their formats are fixed by hardware constraints).", 0),
    ("WOULD DISPUTE: real-building thermal simulation from 2D — which is why we scope it out explicitly.", 0),
])

table("Hypotheses (pre-registered)",
      ["#", "Hypothesis", "Where tested"],
      [
          ["S1/S2", "I_V rises with tier for physical targets; saturation is task-dependent", "main grid"],
          ["S3", "capacity compensation is partial and costs ≥5× data", "capacity/sample analysis"],
          ["S4", "annotation gain ≈ 0 residential, large institutional (the headline if it holds)", "MSD vs ArchCAD regimes"],
          ["S5/S6", "fine-tuned transformers preserve probe ordering; functional > geometric hierarchy", "confirmation phase"],
          ["S7/S8", "code-graph ladder transfers; analytic target certifies the simulated one", "generality + consistency"],
          ["S9", "non-monotone dips for small probe families, vanishing with capacity (Prop. 1)", "grid stratified by V"],
      ])

bullets("Where we are (day 8 of the 8-week plan)", [
    ("DONE: six-tier schema + tested forgetting maps; calibrated instrument (V0–V6, PASS); pipeline validated vs ground truth (room F1 0.986); 2,883-and-growing tier-graph corpus; owner-verified flagship building; door-detector retrain program (11× gain, v2 pending); 4 corpora landed and dissected; methods review folded in.", 0),
    ("RUNNING NOW: corpus completion sweep (patched, 0 new failures); full-corpus re-scoring; door retrain v2 prep.", 0),
    ("NEXT (this week): the stage-gate probing grid on T0–T3 — the experiment this entire infrastructure exists to run.", 0),
    ("Then: gate verdict → manual T4/T5 annotation (timed) → capacity/sample analyses → confirmation phase → write-up.", 0),
])

bullets("Decisions requested this week", [
    ("Grid scale: include free-structure controls (+≈30% cells) — recommended yes (they make C3 defensible).", 0),
    ("Fifth target: add space-syntax integration (zero annotation cost) — recommended yes.", 0),
    ("Bayes-error estimator: hold as rebuttal contingency (recommended) vs. run in main grid.", 0),
    ("Door verdict: on the real sheets, are the retrained detector's extra doors real or hallucinated? (eval_real.png) — shapes retrain v2.", 0),
    ("Annotation-hours planning: if the stage gate opens, ≤40 h of timed T4/T5 annotation in weeks 3–4.", 0),
])

bullets("Everything is inspectable (reproducibility posture)", [
    ("Repo: github.com/YaqoobAnsari/InfoLadder — every number traces to a run manifest (git SHA + config hash + seed) in an append-only registry.", 0),
    ("Browsable corpus: results/corpora/ — per building: source floorplan, tier JSONs, per-tier overlays, report card.", 0),
    ("Phase reports with figures: results/reports/; methods review: docs/audits/; decision log: docs/DECISIONS.md (18 records).", 0),
    ("External accountability: all code and results are independently monitored and graded.", 0),
])

# ============================== RENDERERS ====================================

def _img_size(path):
    with PILImage.open(path) as im:
        return im.size


def render_pptx(out_path: Path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt

    EMU_W, EMU_H = Inches(13.333), Inches(7.5)
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    blank = prs.slide_layouts[6]

    C_INK = RGBColor(0x0B, 0x0B, 0x0B)
    C_SUB = RGBColor(0x52, 0x51, 0x4E)
    C_ACC = RGBColor(0x2A, 0x78, 0xD6)
    C_BG = RGBColor(0xFC, 0xFC, 0xFB)

    def new_slide():
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C_BG
        return s

    def tb(s, x, y, w, h):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        return box.text_frame

    def slide_title(s, text, color=C_INK, size=28):
        tf = tb(s, 0.6, 0.35, 12.1, 0.9)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        bar = s.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(1.6), Emu(28575))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_ACC
        bar.line.fill.background()

    for spec in S:
        k = spec["kind"]
        s = new_slide()
        if k == "title":
            tf = tb(s, 1.0, 2.0, 11.3, 1.2)
            p = tf.paragraphs[0]
            p.text = spec["title"]
            p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = C_ACC
            tf2 = tb(s, 1.0, 3.2, 11.3, 1.6)
            for i, line in enumerate(spec["sub"].split("\n")):
                p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p.text = line; p.font.size = Pt(22); p.font.color.rgb = C_INK
            tf3 = tb(s, 1.0, 5.6, 11.3, 1.4)
            for i, line in enumerate(spec["meta"]):
                p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
                p.text = line; p.font.size = Pt(14); p.font.color.rgb = C_SUB
        elif k == "section":
            tf = tb(s, 1.0, 2.6, 11.3, 2.4)
            p = tf.paragraphs[0]
            p.text = f"{spec['n']}."
            p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = C_ACC
            p2 = tf.add_paragraph()
            p2.text = spec["title"]; p2.font.size = Pt(34); p2.font.bold = True; p2.font.color.rgb = C_INK
            if spec["blurb"]:
                p3 = tf.add_paragraph()
                p3.text = spec["blurb"]; p3.font.size = Pt(16); p3.font.color.rgb = C_SUB
        elif k == "statement":
            slide_title(s, spec["title"])
            tf = tb(s, 1.2, 1.8, 10.9, 5.0)
            for i, line in enumerate(spec["lines"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line; p.font.size = Pt(20); p.font.color.rgb = C_INK
                p.space_after = Pt(10)
        elif k == "bullets":
            slide_title(s, spec["title"])
            tf = tb(s, 0.7, 1.45, 12.0, 5.6)
            first = True
            for text, lvl in spec["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = ("• " if lvl == 0 else "– ") + text
                p.level = lvl
                p.font.size = Pt(16 if lvl == 0 else 14)
                p.font.color.rgb = C_INK if lvl == 0 else C_SUB
                p.space_after = Pt(7)
            if spec["note"]:
                nf = tb(s, 0.7, 6.9, 12.0, 0.5)
                p = nf.paragraphs[0]; p.text = spec["note"]; p.font.size = Pt(11); p.font.color.rgb = C_SUB
        elif k == "table":
            slide_title(s, spec["title"], size=24)
            rows, cols = len(spec["rows"]) + 1, len(spec["headers"])
            widths = [Inches(12.1 / cols)] * cols
            shape = s.shapes.add_table(rows, cols, Inches(0.6), Inches(1.4),
                                       Inches(12.1), Inches(0.4 * rows))
            t = shape.table
            for j, hcell in enumerate(spec["headers"]):
                c = t.cell(0, j); c.text = hcell
                c.text_frame.paragraphs[0].font.size = Pt(12)
                c.text_frame.paragraphs[0].font.bold = True
            for i, row in enumerate(spec["rows"], start=1):
                for j, val in enumerate(row):
                    c = t.cell(i, j); c.text = str(val)
                    for p in c.text_frame.paragraphs:
                        p.font.size = Pt(11)
            if spec["note"]:
                nf = tb(s, 0.6, 6.95, 12.1, 0.45)
                p = nf.paragraphs[0]; p.text = spec["note"]; p.font.size = Pt(10); p.font.color.rgb = C_SUB
        elif k in ("image", "imagebul"):
            slide_title(s, spec["title"], size=24)
            iw, ih = _img_size(spec["img"])
            if k == "image":
                max_w, max_h, x0 = 11.6, 5.0, 0.85
            else:
                max_w, max_h, x0 = 6.9, 5.1, 6.1
            scale = min(max_w / iw, max_h / ih)
            w, h = iw * scale, ih * scale
            s.shapes.add_picture(str(spec["img"]), Inches(x0 + (max_w - w) / 2),
                                 Inches(1.4), Inches(w), Inches(h))
            if k == "imagebul":
                tf = tb(s, 0.6, 1.5, 5.3, 4.9)
                first = True
                for text, lvl in spec["items"]:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = ("• " if lvl == 0 else "– ") + text
                    p.font.size = Pt(14 if lvl == 0 else 12)
                    p.font.color.rgb = C_INK if lvl == 0 else C_SUB
                    p.space_after = Pt(6)
            if spec["caption"]:
                nf = tb(s, 0.6, 6.6, 12.1, 0.8)
                p = nf.paragraphs[0]; p.text = spec["caption"]; p.font.size = Pt(11); p.font.color.rgb = C_SUB
    prs.save(str(out_path))


def render_pdf(out_path: Path):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    def page():
        fig = plt.figure(figsize=(13.333, 7.5))
        fig.patch.set_facecolor(BG)
        return fig

    def head(fig, text, size=22):
        fig.text(0.045, 0.925, text, fontsize=size, fontweight="bold", color=INK, va="top")
        fig.add_artist(plt.Line2D([0.045, 0.17], [0.845, 0.845], color=ACC, lw=3,
                                  transform=fig.transFigure))

    with PdfPages(out_path) as pdf:
        for spec in S:
            k = spec["kind"]
            fig = page()
            if k == "title":
                fig.text(0.08, 0.68, spec["title"], fontsize=44, fontweight="bold", color=ACC)
                fig.text(0.08, 0.52, spec["sub"], fontsize=18, color=INK)
                fig.text(0.08, 0.2, "\n".join(spec["meta"]), fontsize=11, color=SUB)
            elif k == "section":
                fig.text(0.08, 0.60, f"{spec['n']}.", fontsize=52, fontweight="bold", color=ACC)
                fig.text(0.08, 0.48, spec["title"], fontsize=28, fontweight="bold", color=INK)
                if spec["blurb"]:
                    fig.text(0.08, 0.38, spec["blurb"], fontsize=13, color=SUB, wrap=True)
            elif k == "statement":
                head(fig, spec["title"])
                fig.text(0.08, 0.72, "\n\n".join(spec["lines"]), fontsize=15, color=INK,
                         va="top", wrap=True)
            elif k == "bullets":
                head(fig, spec["title"])
                y = 0.78
                for text, lvl in spec["items"]:
                    x = 0.06 + 0.03 * lvl
                    mark = "• " if lvl == 0 else "– "
                    import textwrap
                    wrapped = textwrap.fill(mark + text, width=int(118 - 6 * lvl))
                    n_lines = wrapped.count("\n") + 1
                    fig.text(x, y, wrapped, fontsize=12.5 if lvl == 0 else 11,
                             color=INK if lvl == 0 else SUB, va="top")
                    y -= 0.052 * n_lines + 0.012
                if spec["note"]:
                    fig.text(0.06, 0.045, spec["note"], fontsize=9, color=SUB)
            elif k == "table":
                head(fig, spec["title"], size=19)
                ax = fig.add_axes([0.04, 0.06, 0.92, 0.72]); ax.axis("off")
                tbl = ax.table(cellText=[[str(v) for v in r] for r in spec["rows"]],
                               colLabels=spec["headers"], loc="upper center",
                               cellLoc="left", colLoc="left")
                tbl.auto_set_font_size(False)
                tbl.set_fontsize(9.5)
                tbl.scale(1, 2.1)
                for (r, c), cell in tbl.get_celld().items():
                    cell.set_edgecolor("#dddbd4")
                    cell.set_facecolor("#ffffff" if r else "#eef3fb")
                    if r == 0:
                        cell.set_text_props(fontweight="bold")
                if spec["note"]:
                    fig.text(0.05, 0.03, spec["note"], fontsize=9, color=SUB)
            elif k in ("image", "imagebul"):
                head(fig, spec["title"], size=19)
                img = plt.imread(str(spec["img"]))
                if k == "image":
                    ax = fig.add_axes([0.08, 0.14, 0.84, 0.68])
                else:
                    ax = fig.add_axes([0.47, 0.13, 0.5, 0.68])
                    y = 0.76
                    import textwrap
                    for text, lvl in spec["items"]:
                        mark = "• " if lvl == 0 else "– "
                        wrapped = textwrap.fill(mark + text, width=52)
                        n_lines = wrapped.count("\n") + 1
                        fig.text(0.05 + 0.02 * lvl, y, wrapped,
                                 fontsize=11 if lvl == 0 else 10,
                                 color=INK if lvl == 0 else SUB, va="top")
                        y -= 0.045 * n_lines + 0.012
                ax.imshow(img)
                ax.axis("off")
                if spec["caption"]:
                    import textwrap
                    fig.text(0.06, 0.055, textwrap.fill(spec["caption"], width=150),
                             fontsize=9.5, color=SUB)
            pdf.savefig(fig)
            plt.close(fig)


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    missing = [k for k, p in IMG.items() if not p.exists()]
    if missing:
        raise SystemExit(f"missing image assets: {missing}")
    render_pptx(OUT / f"{STEM}.pptx")
    print(f"wrote {OUT / (STEM + '.pptx')}")
    render_pdf(OUT / f"{STEM}.pdf")
    print(f"wrote {OUT / (STEM + '.pdf')}")
    print(f"slides: {len(S)}")
